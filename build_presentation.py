# -*- coding: utf-8 -*-
import os
import sys
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

# -----------------------------------------------------------------------------
# Color Palette Definitions
# -----------------------------------------------------------------------------
NAVY_DARK = RGBColor(11, 19, 43)      # #0B132B (Base for dark slides)
NAVY_MEDIUM = RGBColor(28, 37, 65)    # #1C2541 (Headers, dark cards)
BG_LIGHT = RGBColor(248, 250, 252)    # #F8FAFC (Content slide background)
CYAN_ACCENT = RGBColor(0, 180, 216)   # #00B4D8 (Highlight accents, step numbers)
CORAL_WARNING = RGBColor(255, 90, 95) # #FF5A5F (Alerts, times, warning boxes)
WHITE = RGBColor(255, 255, 255)
GRAY_TEXT = RGBColor(74, 85, 104)     # #4A5568
DARK_TEXT = RGBColor(28, 37, 65)

# -----------------------------------------------------------------------------
# Helper Functions for Transparency and Formatting
# -----------------------------------------------------------------------------
def set_shape_transparency(shape, opacity):
    """
    Sets the fill transparency of a shape using OpenXML manipulation.
    opacity: float between 0.0 (transparent) and 1.0 (opaque)
    """
    try:
        from pptx.oxml.xmlchemy import OxmlElement
        shape.fill.solid()
        solidFill = shape.element.spPr.find('{http://schemas.openxmlformats.org/drawingml/2006/main}solidFill')
        if solidFill is not None:
            # Check if alpha already exists
            alpha = solidFill.find('{http://schemas.openxmlformats.org/drawingml/2006/main}alpha')
            if alpha is not None:
                alpha.set('val', str(int(opacity * 100000)))
            else:
                alpha = OxmlElement('a:alpha')
                alpha.set('val', str(int(opacity * 100000)))
                solidFill.append(alpha)
    except Exception as e:
        print("Failed to set shape transparency:", e)

def map_icon_to_text(text):
    """
    Automatically maps vector-style emojis next to text points.
    """
    text_lower = text.lower()
    
    # Clean leadings
    cleaned = text.strip()
    if cleaned.startswith(("•", "-", "*", "🔹", "📌", "✔")):
        cleaned = cleaned[1:].strip()
        
    if cleaned.startswith(("🎯", "📋", "⚙️", "⏱️", "📈", "💡", "🎤", "❓", "👁️", "🚣", "🗣️", "🔹", "🔍", "🌪️", "⛵", "🦺", "🛡️", "🚨", "⭐", "📚", "🛠️", "📢", "🥤", "🪵", "🪨", "🌿")):
        return cleaned
        
    if "chuột rút" in text_lower:
        return "🔍 " + cleaned
    elif "nước xoáy" in text_lower or "dòng nước xoáy" in text_lower:
        return "🌪️ " + cleaned
    elif "thuyền" in text_lower or "lật thuyền" in text_lower:
        return "⛵ " + cleaned
    elif "áo phao" in text_lower:
        return "🦺 " + cleaned
    elif "khởi động" in text_lower:
        return "🛡️ " + cleaned
    elif "cảnh báo" in text_lower or "nguy cơ" in text_lower or "nguy hiểm" in text_lower:
        return "🚨 " + cleaned
    elif "bất ngờ" in text_lower:
        return "⭐ " + cleaned
    elif "kiến thức" in text_lower:
        return "📚 " + cleaned
    elif "kỹ năng" in text_lower:
        return "🛠️ " + cleaned
    elif "thông điệp" in text_lower:
        return "📢 " + cleaned
    elif "cam kết" in text_lower or "tôi sẽ" in text_lower or "em sẽ" in text_lower:
        return "👁️ " + cleaned
    elif "can nhựa" in text_lower or "can rỗng" in text_lower:
        return "🥤 " + cleaned
    elif "khúc cây" in text_lower:
        return "🪵 " + cleaned
    elif "đá" in text_lower or "cục đá" in text_lower:
        return "🪨 " + cleaned
    elif "cành cây" in text_lower:
        return "🌿 " + cleaned
    elif "phóng viên" in text_lower or "nova news" in text_lower:
        return "🎤 " + cleaned
    elif "mục tiêu" in text_lower:
        return "🎯 " + cleaned
    elif "kết quả" in text_lower:
        return "📈 " + cleaned
    else:
        return "🔹 " + cleaned

def add_textbox(slide, left, top, width, height, text, font_name="Arial", font_size=14, font_color=WHITE, bold=False, align=PP_ALIGN.LEFT, word_wrap=True):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = word_wrap
    tf.margin_left = tf.margin_top = tf.margin_right = tf.margin_bottom = 0
    p = tf.paragraphs[0]
    p.text = text
    p.font.name = font_name
    p.font.size = Pt(font_size)
    p.font.color.rgb = font_color
    p.font.bold = bold
    p.alignment = align
    return txBox

def add_bullet_points(slide, left, top, width, height, bullets, font_name="Calibri", font_size=14, font_color=DARK_TEXT):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_top = tf.margin_right = tf.margin_bottom = 0
    for idx, b in enumerate(bullets):
        if idx == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        # Map icon to text automatically
        p.text = map_icon_to_text(b)
        p.font.name = font_name
        p.font.size = Pt(font_size)
        p.font.color.rgb = font_color
        p.space_after = Pt(6)
    return txBox

def add_native_table(slide, left, top, width, height, headers, rows_data):
    """
    Inserts an editable native PowerPoint table.
    """
    rows = len(rows_data) + 1
    cols = len(headers)
    table_shape = slide.shapes.add_table(rows, cols, left, top, width, height)
    table = table_shape.table
    
    # Format headers
    for c_idx, h_text in enumerate(headers):
        cell = table.cell(0, c_idx)
        cell.text = h_text
        cell.fill.solid()
        cell.fill.fore_color.rgb = NAVY_MEDIUM
        p = cell.text_frame.paragraphs[0]
        p.font.name = "Arial"
        p.font.size = Pt(12)
        p.font.bold = True
        p.font.color.rgb = WHITE
        p.alignment = PP_ALIGN.CENTER
        
    # Format cells
    for r_idx, row_data in enumerate(rows_data):
        for c_idx, val in enumerate(row_data):
            cell = table.cell(r_idx + 1, c_idx)
            cell.text = val
            cell.fill.solid()
            cell.fill.fore_color.rgb = WHITE
            p = cell.text_frame.paragraphs[0]
            p.font.name = "Arial"
            p.font.size = Pt(11)
            p.font.color.rgb = DARK_TEXT
            p.alignment = PP_ALIGN.LEFT
            
    return table_shape

# -----------------------------------------------------------------------------
# Background Setup Functions (Blurred Newsroom + Opacity Color Overlay)
# -----------------------------------------------------------------------------
def set_dark_background(slide):
    # Base blurred dark newsroom image
    bg_img = "images/newsroom_dark_bg.png"
    if os.path.exists(bg_img):
        slide.shapes.add_picture(bg_img, 0, 0, width=Inches(13.333), height=Inches(7.5))
        
    # Semi-transparent dark navy overlay
    overlay = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.333), Inches(7.5))
    overlay.fill.solid()
    overlay.fill.fore_color.rgb = NAVY_DARK
    overlay.line.color.rgb = NAVY_DARK
    set_shape_transparency(overlay, 0.85)

def set_light_background(slide, stage_num=None, stage_name=None):
    # Base blurred light newsroom image
    bg_img = "images/newsroom_light_bg.png"
    if os.path.exists(bg_img):
        slide.shapes.add_picture(bg_img, 0, 0, width=Inches(13.333), height=Inches(7.5))
        
    # Semi-transparent light gray/slate overlay
    overlay = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.333), Inches(7.5))
    overlay.fill.solid()
    overlay.fill.fore_color.rgb = BG_LIGHT
    overlay.line.color.rgb = BG_LIGHT
    set_shape_transparency(overlay, 0.93)
    
    # Top banner (Navy, slim)
    banner = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.333), Inches(0.8))
    banner.fill.solid()
    banner.fill.fore_color.rgb = NAVY_MEDIUM
    banner.line.color.rgb = NAVY_MEDIUM
    
    # Lesson Title in top banner
    add_textbox(slide, Inches(0.5), Inches(0.2), Inches(6.0), Inches(0.4), "AN TOÀN CÁ NHÂN • PHÒNG TRÁNH ĐUỐI NƯỚC", font_name="Arial", font_size=11, font_color=CYAN_ACCENT, bold=True)
    
    # Stage indicator on the right of top banner
    if stage_num and stage_name:
        add_textbox(slide, Inches(7.0), Inches(0.2), Inches(5.8), Inches(0.4), f"Giai đoạn {stage_num}: {stage_name}", font_name="Arial", font_size=11, font_color=WHITE, bold=True, align=PP_ALIGN.RIGHT)

# -----------------------------------------------------------------------------
# Main Template Generator Functions
# -----------------------------------------------------------------------------
def create_stage_intro_slide(prs, stage_num, stage_name, dna_type):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_dark_background(slide)
    
    # Large DNA Label
    add_textbox(slide, Inches(1.0), Inches(2.0), Inches(11.333), Inches(0.5), f"LESSON DNA: {dna_type.upper()}", font_name="Arial", font_size=16, font_color=CYAN_ACCENT, bold=True)
    
    # Stage Number
    add_textbox(slide, Inches(1.0), Inches(2.6), Inches(11.333), Inches(0.8), f"GIAI ĐOẠN 0{stage_num}", font_name="Arial", font_size=36, font_color=CORAL_WARNING, bold=True)
    
    # Stage Name
    add_textbox(slide, Inches(1.0), Inches(3.5), Inches(11.333), Inches(1.5), stage_name.upper(), font_name="Arial", font_size=44, font_color=WHITE, bold=True)
    
    # Add a thin line under the stage name
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(1.0), Inches(5.2), Inches(4.0), Inches(0.04))
    line.fill.solid()
    line.fill.fore_color.rgb = CYAN_ACCENT
    line.line.color.rgb = CYAN_ACCENT
    
    # Role description at bottom
    add_textbox(slide, Inches(1.0), Inches(5.5), Inches(11.333), Inches(0.5), "🎤 Nhiệm vụ Phóng viên Nova News", font_name="Calibri", font_size=18, font_color=RGBColor(200, 200, 200))
    return slide

def create_act01_slide(prs, stage_num, stage_name, act_title, goal, image_path, extra_desc=None, table_info=None, notebook_info=None, journal_info=None, commitment_info=None):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_light_background(slide, stage_num=stage_num, stage_name=stage_name)
    
    # Activity Title
    add_textbox(slide, Inches(0.8), Inches(1.1), Inches(6.0), Inches(0.6), f"Hoạt động: {act_title.upper()}", font_name="Arial", font_size=22, font_color=DARK_TEXT, bold=True)
    
    # Card for Goal & Info
    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(1.8), Inches(5.8), Inches(4.8))
    card.fill.solid()
    card.fill.fore_color.rgb = WHITE
    card.line.color.rgb = RGBColor(226, 232, 240)
    
    # Card Title: Mục tiêu
    add_textbox(slide, Inches(1.2), Inches(2.1), Inches(5.0), Inches(0.4), "🎯 MỤC TIÊU HOẠT ĐỘNG", font_name="Arial", font_size=13, font_color=CYAN_ACCENT, bold=True)
    
    # Goal bullets
    goals = [g.strip() for g in goal.split("\n") if g.strip()]
    goals = [g[1:].strip() if g.startswith("•") or g.startswith("-") else g for g in goals]
    add_bullet_points(slide, Inches(1.2), Inches(2.6), Inches(5.0), Inches(1.6), goals, font_size=15, font_color=DARK_TEXT)
    
    # Extra description if available
    if extra_desc:
        desc_card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1.1), Inches(4.3), Inches(5.2), Inches(2.0))
        desc_card.fill.solid()
        desc_card.fill.fore_color.rgb = BG_LIGHT
        desc_card.line.color.rgb = RGBColor(226, 232, 240)
        add_textbox(slide, Inches(1.3), Inches(4.5), Inches(4.8), Inches(1.6), extra_desc, font_name="Calibri", font_size=14, font_color=GRAY_TEXT)
    
    # Image on the right
    if image_path and os.path.exists(image_path):
        slide.shapes.add_picture(image_path, Inches(7.0), Inches(1.8), width=Inches(5.5), height=Inches(4.8))
    else:
        pl = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(7.0), Inches(1.8), Inches(5.5), Inches(4.8))
        pl.fill.solid()
        pl.fill.fore_color.rgb = NAVY_MEDIUM
        pl.line.color.rgb = CYAN_ACCENT
        add_textbox(slide, Inches(7.2), Inches(3.8), Inches(5.1), Inches(1.0), f"[Visual Illustration]\n{image_path}", font_name="Arial", font_size=16, bold=True, align=PP_ALIGN.CENTER)
        
    # --- Overlay structures if required (Table, Notebook, Reflection, Commitment) ---
    if table_info:
        headers, rows_data = table_info
        add_native_table(slide, Inches(7.5), Inches(2.6), Inches(4.5), Inches(3.2), headers, rows_data)
        
    elif notebook_info:
        left_text, right_text = notebook_info
        # Left page
        add_textbox(slide, Inches(7.35), Inches(2.3), Inches(2.2), Inches(3.8), left_text, font_name="Arial", font_size=10, font_color=DARK_TEXT, bold=True)
        # Right page
        add_textbox(slide, Inches(9.85), Inches(2.3), Inches(2.2), Inches(3.8), right_text, font_name="Arial", font_size=10, font_color=DARK_TEXT, bold=True)
        
    elif journal_info:
        add_textbox(slide, Inches(7.45), Inches(3.2), Inches(4.6), Inches(3.0), journal_info, font_name="Arial", font_size=12, font_color=NAVY_MEDIUM, bold=True)
        
    elif commitment_info:
        add_textbox(slide, Inches(7.8), Inches(3.5), Inches(3.9), Inches(2.2), commitment_info, font_name="Arial", font_size=12, font_color=NAVY_MEDIUM, bold=True, align=PP_ALIGN.CENTER)
        
    return slide

def create_act02_slide(prs, stage_num, stage_name, act_title, steps):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_light_background(slide, stage_num=stage_num, stage_name=stage_name)
    
    # Activity Title
    add_textbox(slide, Inches(0.8), Inches(1.1), Inches(11.7), Inches(0.6), f"Hoạt động: {act_title.upper()} • HƯỚNG DẪN THỰC HIỆN", font_name="Arial", font_size=20, font_color=DARK_TEXT, bold=True)
    
    n_steps = len(steps)
    if n_steps == 1:
        card_w = Inches(8.0)
        card_h = Inches(3.5)
        left = Inches(2.666)
        top = Inches(2.5)
        
        card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, card_w, card_h)
        card.fill.solid()
        card.fill.fore_color.rgb = WHITE
        card.line.color.rgb = CYAN_ACCENT
        card.line.width = Pt(2)
        
        num_circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, left + Inches(0.6), top + Inches(0.6), Inches(1.0), Inches(1.0))
        num_circle.fill.solid()
        num_circle.fill.fore_color.rgb = CYAN_ACCENT
        num_circle.line.color.rgb = CYAN_ACCENT
        add_textbox(slide, left + Inches(0.6), top + Inches(0.85), Inches(1.0), Inches(0.5), "1", font_name="Arial", font_size=28, font_color=WHITE, bold=True, align=PP_ALIGN.CENTER)
        
        add_textbox(slide, left + Inches(2.0), top + Inches(0.6), card_w - Inches(2.6), Inches(2.3), "👉 " + steps[0], font_name="Arial", font_size=22, font_color=DARK_TEXT, bold=True)
        
    elif n_steps == 2:
        card_w = Inches(5.4)
        card_h = Inches(3.8)
        spacing = Inches(0.9)
        start_left = Inches(0.8)
        top = Inches(2.3)
        
        for i, step in enumerate(steps):
            left = start_left + i * (card_w + spacing)
            card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, card_w, card_h)
            card.fill.solid()
            card.fill.fore_color.rgb = WHITE
            card.line.color.rgb = CYAN_ACCENT
            card.line.width = Pt(2)
            
            num_circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, left + Inches(0.5), top + Inches(0.5), Inches(0.9), Inches(0.9))
            num_circle.fill.solid()
            num_circle.fill.fore_color.rgb = CYAN_ACCENT
            num_circle.line.color.rgb = CYAN_ACCENT
            add_textbox(slide, left + Inches(0.5), top + Inches(0.7), Inches(0.9), Inches(0.5), str(i+1), font_name="Arial", font_size=24, font_color=WHITE, bold=True, align=PP_ALIGN.CENTER)
            
            add_textbox(slide, left + Inches(1.7), top + Inches(0.5), card_w - Inches(2.1), card_h - Inches(1.0), "👉 " + step, font_name="Arial", font_size=20, font_color=DARK_TEXT, bold=True)
            
    elif n_steps >= 3:
        card_w = Inches(3.6)
        card_h = Inches(4.2)
        spacing = Inches(0.5)
        start_left = Inches(0.8)
        top = Inches(2.2)
        
        for i, step in enumerate(steps[:3]):
            left = start_left + i * (card_w + spacing)
            card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, card_w, card_h)
            card.fill.solid()
            card.fill.fore_color.rgb = WHITE
            card.line.color.rgb = CYAN_ACCENT
            card.line.width = Pt(2)
            
            num_circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, left + Inches(0.4), top + Inches(0.4), Inches(0.8), Inches(0.8))
            num_circle.fill.solid()
            num_circle.fill.fore_color.rgb = CYAN_ACCENT
            num_circle.line.color.rgb = CYAN_ACCENT
            add_textbox(slide, left + Inches(0.4), top + Inches(0.55), Inches(0.8), Inches(0.5), str(i+1), font_name="Arial", font_size=20, font_color=WHITE, bold=True, align=PP_ALIGN.CENTER)
            
            add_textbox(slide, left + Inches(0.4), top + Inches(1.4), card_w - Inches(0.8), card_h - Inches(1.6), "👉 " + step, font_name="Arial", font_size=18, font_color=DARK_TEXT, bold=True)
            
    return slide

def create_act03_slide(prs, stage_num, stage_name, act_title, prompt_text, product, format_org, time_mins):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_light_background(slide, stage_num=stage_num, stage_name=stage_name)
    
    # Activity Title
    add_textbox(slide, Inches(0.8), Inches(1.1), Inches(11.7), Inches(0.6), f"Hoạt động: {act_title.upper()} • KHÔNG GIAN HOẠT ĐỘNG", font_name="Arial", font_size=20, font_color=DARK_TEXT, bold=True)
    
    # Left Card: Đề bài
    left_card_w = Inches(6.8)
    left_card_h = Inches(4.8)
    left_card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(1.8), left_card_w, left_card_h)
    left_card.fill.solid()
    left_card.fill.fore_color.rgb = WHITE
    left_card.line.color.rgb = RGBColor(226, 232, 240)
    
    header_left = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(1.8), left_card_w, Inches(0.8))
    header_left.fill.solid()
    header_left.fill.fore_color.rgb = NAVY_MEDIUM
    header_left.line.color.rgb = NAVY_MEDIUM
    add_textbox(slide, Inches(1.2), Inches(2.0), left_card_w - Inches(0.8), Inches(0.4), "📋 ĐỀ BÀI THẢO LUẬN", font_name="Arial", font_size=15, font_color=WHITE, bold=True)
    
    # Đề bài text with icons
    prompt_lines = [map_icon_to_text(l) for l in prompt_text.split("\n") if l.strip()]
    add_bullet_points(slide, Inches(1.2), Inches(2.9), left_card_w - Inches(0.8), left_card_h - Inches(1.3), prompt_lines, font_size=15)
    
    # Right Card: Thông tin tổ chức
    right_card_w = Inches(4.4)
    right_card_h = Inches(4.8)
    right_card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(8.1), Inches(1.8), right_card_w, right_card_h)
    right_card.fill.solid()
    right_card.fill.fore_color.rgb = WHITE
    right_card.line.color.rgb = RGBColor(226, 232, 240)
    
    header_right = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(8.1), Inches(1.8), right_card_w, Inches(0.8))
    header_right.fill.solid()
    header_right.fill.fore_color.rgb = CYAN_ACCENT
    header_right.line.color.rgb = CYAN_ACCENT
    add_textbox(slide, Inches(8.5), Inches(2.0), right_card_w - Inches(0.8), Inches(0.4), "⚙️ YÊU CẦU HOÀN THÀNH", font_name="Arial", font_size=15, font_color=WHITE, bold=True)
    
    # Content on right card
    add_textbox(slide, Inches(8.5), Inches(2.8), right_card_w - Inches(0.8), Inches(0.3), "Sản phẩm cần hoàn thành:", font_name="Arial", font_size=13, font_color=GRAY_TEXT, bold=True)
    add_textbox(slide, Inches(8.5), Inches(3.1), right_card_w - Inches(0.8), Inches(0.8), "📝 " + product, font_name="Calibri", font_size=15, font_color=DARK_TEXT)
    
    add_textbox(slide, Inches(8.5), Inches(4.1), right_card_w - Inches(0.8), Inches(0.3), "Hình thức tổ chức:", font_name="Arial", font_size=13, font_color=GRAY_TEXT, bold=True)
    add_textbox(slide, Inches(8.5), Inches(4.4), right_card_w - Inches(0.8), Inches(0.4), "👥 " + format_org, font_name="Calibri", font_size=16, font_color=DARK_TEXT)
    
    # Time Badge
    badge_w = Inches(2.2)
    badge_h = Inches(0.6)
    badge = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(8.5), Inches(5.5), badge_w, badge_h)
    badge.fill.solid()
    badge.fill.fore_color.rgb = CORAL_WARNING
    badge.line.color.rgb = CORAL_WARNING
    add_textbox(slide, Inches(8.5), Inches(5.6), badge_w, Inches(0.4), f"⏱️ {time_mins}", font_name="Arial", font_size=15, font_color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    
    return slide

def create_act04_slide(prs, stage_num, stage_name, act_title, results, key_message):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_light_background(slide, stage_num=stage_num, stage_name=stage_name)
    
    # Activity Title
    add_textbox(slide, Inches(0.8), Inches(1.1), Inches(11.7), Inches(0.6), f"Hoạt động: {act_title.upper()} • BÁO CÁO & CHỐT KIẾN THỨC", font_name="Arial", font_size=20, font_color=DARK_TEXT, bold=True)
    
    # Left Card: Kết quả
    left_card_w = Inches(6.8)
    left_card_h = Inches(4.8)
    left_card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(1.8), left_card_w, left_card_h)
    left_card.fill.solid()
    left_card.fill.fore_color.rgb = WHITE
    left_card.line.color.rgb = RGBColor(226, 232, 240)
    
    header_left = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(1.8), left_card_w, Inches(0.8))
    header_left.fill.solid()
    header_left.fill.fore_color.rgb = NAVY_MEDIUM
    header_left.line.color.rgb = NAVY_MEDIUM
    add_textbox(slide, Inches(1.2), Inches(2.0), left_card_w - Inches(0.8), Inches(0.4), "📈 KẾT QUẢ THU HOẠCH", font_name="Arial", font_size=15, font_color=WHITE, bold=True)
    
    # Bullets for results
    results_list = [r.strip() for r in results.split("\n") if r.strip()]
    results_list = [r[1:].strip() if r.startswith("•") or r.startswith("-") else r for r in results_list]
    add_bullet_points(slide, Inches(1.2), Inches(2.8), left_card_w - Inches(0.8), left_card_h - Inches(1.2), results_list, font_size=16, font_color=DARK_TEXT)
    
    # Right Card: Thông điệp chính (Callout Box)
    right_card_w = Inches(4.4)
    right_card_h = Inches(4.8)
    
    right_card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(8.1), Inches(1.8), right_card_w, right_card_h)
    right_card.fill.solid()
    right_card.fill.fore_color.rgb = RGBColor(224, 251, 252) # Light aqua background
    right_card.line.color.rgb = CYAN_ACCENT
    right_card.line.width = Pt(2.5)
    
    add_textbox(slide, Inches(8.5), Inches(2.2), right_card_w - Inches(0.8), Inches(0.4), "💡 THÔNG ĐIỆP CHÍNH", font_name="Arial", font_size=16, font_color=CYAN_ACCENT, bold=True)
    
    add_textbox(slide, Inches(8.5), Inches(2.9), right_card_w - Inches(0.8), Inches(3.3), "📢 " + key_message, font_name="Arial", font_size=18, font_color=DARK_TEXT, bold=True)
    
    return slide

def create_act05_slide(prs, stage_num, stage_name, act_title, explored_text, next_question):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_dark_background(slide)
    
    # Top header
    add_textbox(slide, Inches(1.0), Inches(1.0), Inches(11.333), Inches(0.5), f"HOÀN THÀNH HOẠT ĐỘNG: {act_title.upper()}", font_name="Arial", font_size=14, font_color=CYAN_ACCENT, bold=True)
    
    # Main Exploration Box
    add_textbox(slide, Inches(1.0), Inches(1.8), Inches(11.333), Inches(0.4), "Chúng ta vừa khám phá:", font_name="Calibri", font_size=18, font_color=RGBColor(160, 174, 192))
    add_textbox(slide, Inches(1.0), Inches(2.3), Inches(11.333), Inches(1.2), "👉 " + explored_text, font_name="Arial", font_size=28, font_color=WHITE, bold=True)
    
    # A thick divider line
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(1.0), Inches(3.8), Inches(11.333), Inches(0.03))
    line.fill.solid()
    line.fill.fore_color.rgb = NAVY_MEDIUM
    line.line.color.rgb = NAVY_MEDIUM
    
    # Next Question
    add_textbox(slide, Inches(1.0), Inches(4.2), Inches(11.333), Inches(0.4), "Câu hỏi tiếp theo dành cho phóng viên:", font_name="Calibri", font_size=18, font_color=CORAL_WARNING)
    add_textbox(slide, Inches(1.0), Inches(4.7), Inches(11.333), Inches(1.5), "❓ " + next_question, font_name="Arial", font_size=32, font_color=CYAN_ACCENT, bold=True)
    
    # Go to next stage button representation
    btn = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(10.0), Inches(6.2), Inches(2.3), Inches(0.6))
    btn.fill.solid()
    btn.fill.fore_color.rgb = CYAN_ACCENT
    btn.line.color.rgb = CYAN_ACCENT
    add_textbox(slide, Inches(10.0), Inches(6.32), Inches(2.3), Inches(0.4), "TIẾP TỤC ➔", font_name="Arial", font_size=14, font_color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    
    return slide

# -----------------------------------------------------------------------------
# Specific Custom Slides
# -----------------------------------------------------------------------------
def create_welcome_slide(prs, image_path):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_dark_background(slide)
    
    # Left Panel
    add_textbox(slide, Inches(0.8), Inches(1.0), Inches(6.0), Inches(0.4), "CHỦ ĐỀ: AN TOÀN CÁ NHÂN", font_name="Arial", font_size=14, font_color=CYAN_ACCENT, bold=True)
    add_textbox(slide, Inches(0.8), Inches(1.5), Inches(6.0), Inches(1.2), "PHÒNG TRÁNH\nĐUỐI NƯỚC", font_name="Arial", font_size=44, font_color=WHITE, bold=True)
    
    role_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(2.9), Inches(3.2), Inches(0.6))
    role_box.fill.solid()
    role_box.fill.fore_color.rgb = CORAL_WARNING
    role_box.line.color.rgb = CORAL_WARNING
    add_textbox(slide, Inches(0.8), Inches(3.05), Inches(3.2), Inches(0.4), "🎤 VAI TRÒ HỌC SINH: PHÓNG VIÊN", font_name="Arial", font_size=13, font_color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    
    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(3.8), Inches(6.0), Inches(3.0))
    card.fill.solid()
    card.fill.fore_color.rgb = NAVY_MEDIUM
    card.line.color.rgb = NAVY_MEDIUM
    
    add_textbox(slide, Inches(1.1), Inches(4.0), Inches(5.4), Inches(0.3), "🎯 MỤC TIÊU BÀI HỌC", font_name="Arial", font_size=13, font_color=CYAN_ACCENT, bold=True)
    
    objectives = [
        "Kiến thức: Nhận biết các nguy cơ đuối nước phổ biến; cách phòng tránh; biết các vật hỗ trợ nổi.",
        "Kỹ năng: Thu thập thông tin; phân tích tình huống; kết nối nguyên nhân và giải pháp.",
        "Thái độ: Có ý thức bảo vệ bản thân và quan tâm đến sự an toàn của người khác."
    ]
    add_bullet_points(slide, Inches(1.1), Inches(4.4), Inches(5.4), Inches(2.2), objectives, font_name="Calibri", font_size=13, font_color=WHITE)
    
    if os.path.exists(image_path):
        slide.shapes.add_picture(image_path, Inches(7.2), Inches(1.0), width=Inches(5.3), height=Inches(5.5))
    else:
        pl = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(7.2), Inches(1.0), Inches(5.3), Inches(5.5))
        pl.fill.solid()
        pl.fill.fore_color.rgb = NAVY_MEDIUM
        pl.line.color.rgb = CYAN_ACCENT
        add_textbox(slide, Inches(7.4), Inches(3.5), Inches(4.9), Inches(1.0), f"[Visual Welcome Image]\n{image_path}", font_name="Arial", font_size=16, bold=True, align=PP_ALIGN.CENTER)

def create_big_question_slide(prs, image_path):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_dark_background(slide)
    
    add_textbox(slide, Inches(1.0), Inches(0.8), Inches(11.333), Inches(0.4), "CÂU HỎI LỚN (BIG QUESTION)", font_name="Arial", font_size=16, font_color=CORAL_WARNING, bold=True, align=PP_ALIGN.CENTER)
    
    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1.5), Inches(1.4), Inches(10.333), Inches(1.2))
    card.fill.solid()
    card.fill.fore_color.rgb = NAVY_MEDIUM
    card.line.color.rgb = CYAN_ACCENT
    card.line.width = Pt(2)
    
    add_textbox(slide, Inches(1.5), Inches(1.7), Inches(10.333), Inches(0.8), "Làm thế nào để phòng tránh đuối nước và giúp người khác an toàn hơn?", font_name="Arial", font_size=22, font_color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    
    if os.path.exists(image_path):
        slide.shapes.add_picture(image_path, Inches(4.5), Inches(2.9), width=Inches(4.3), height=Inches(3.8))
    else:
        pl = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(4.5), Inches(2.9), Inches(4.3), Inches(3.8))
        pl.fill.solid()
        pl.fill.fore_color.rgb = NAVY_MEDIUM
        pl.line.color.rgb = CYAN_ACCENT
        add_textbox(slide, Inches(4.7), Inches(4.5), Inches(3.9), Inches(1.0), f"[Visual Situation]\n{image_path}", font_name="Arial", font_size=14, bold=True, align=PP_ALIGN.CENTER)
        
    desc = "🚨 TÌNH HUỐNG GỢI SUY NGHĨ:\nMột nhóm học sinh đi chơi gần sông, một bạn bị trượt chân xuống nước, các bạn còn lại chưa biết làm gì."
    add_textbox(slide, Inches(9.1), Inches(3.5), Inches(3.4), Inches(2.5), desc, font_name="Calibri", font_size=15, font_color=RGBColor(160, 174, 192))

def create_role_intro_slide(prs, image_path):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_dark_background(slide)
    
    add_textbox(slide, Inches(0.8), Inches(0.8), Inches(11.7), Inches(0.4), "GIỚI THIỆU VAI TRÒ (ROLE INTRODUCTION)", font_name="Arial", font_size=16, font_color=CYAN_ACCENT, bold=True)
    
    add_textbox(slide, Inches(0.8), Inches(1.3), Inches(6.0), Inches(0.6), "VAI TRÒ: PHÓNG VIÊN NOVA NEWS", font_name="Arial", font_size=28, font_color=WHITE, bold=True)
    
    desc = "Các em sẽ đóng vai phóng viên Nova News để điều tra nguy cơ đuối nước và thu thập thông tin giúp trẻ em an toàn hơn."
    add_textbox(slide, Inches(0.8), Inches(2.1), Inches(6.0), Inches(1.0), "🎤 " + desc, font_name="Calibri", font_size=16, font_color=RGBColor(200, 200, 200))
    
    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(3.3), Inches(6.0), Inches(2.4))
    card.fill.solid()
    card.fill.fore_color.rgb = NAVY_MEDIUM
    card.line.color.rgb = NAVY_MEDIUM
    
    add_textbox(slide, Inches(1.1), Inches(3.5), Inches(5.4), Inches(0.3), "📋 NHIỆM VỤ CỦA PHÓNG VIÊN:", font_name="Arial", font_size=13, font_color=CORAL_WARNING, bold=True)
    
    tasks = [
        "Điều tra các nguy cơ đuối nước phổ biến",
        "Tìm cách phòng tránh đuối nước hiệu quả nhất",
        "Thu thập thông tin lập sổ tay giúp trẻ em an toàn hơn"
    ]
    add_bullet_points(slide, Inches(1.1), Inches(3.9), Inches(5.4), Inches(1.6), tasks, font_size=14, font_color=WHITE)
    
    q_card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(5.9), Inches(11.7), Inches(0.8))
    q_card.fill.solid()
    q_card.fill.fore_color.rgb = CORAL_WARNING
    q_card.line.color.rgb = CORAL_WARNING
    add_textbox(slide, Inches(1.1), Inches(6.12), Inches(11.1), Inches(0.4), "❓ CÂU HỎI TRỌNG TÂM: Ai liên quan? Có những góc nhìn nào?", font_name="Arial", font_size=16, font_color=WHITE, bold=True)
    
    if os.path.exists(image_path):
        slide.shapes.add_picture(image_path, Inches(7.2), Inches(1.3), width=Inches(5.3), height=Inches(4.4))
    else:
        pl = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(7.2), Inches(1.3), Inches(5.3), Inches(4.4))
        pl.fill.solid()
        pl.fill.fore_color.rgb = NAVY_MEDIUM
        pl.line.color.rgb = CYAN_ACCENT
        add_textbox(slide, Inches(7.4), Inches(3.2), Inches(4.9), Inches(1.0), f"[Visual Reporter Theme]\n{image_path}", font_name="Arial", font_size=14, bold=True, align=PP_ALIGN.CENTER)

def create_end01_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_dark_background(slide)
    
    add_textbox(slide, Inches(0.8), Inches(0.8), Inches(11.7), Inches(0.5), "TỔNG KẾT BÀI HỌC (KEY TAKEAWAYS)", font_name="Arial", font_size=24, font_color=WHITE, bold=True)
    
    card_w = Inches(3.6)
    card_h = Inches(4.5)
    spacing = Inches(0.45)
    top = Inches(1.8)
    
    # Card 1: Knowledge
    card1 = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), top, card_w, card_h)
    card1.fill.solid()
    card1.fill.fore_color.rgb = NAVY_MEDIUM
    card1.line.color.rgb = CYAN_ACCENT
    card1.line.width = Pt(1.5)
    add_textbox(slide, Inches(1.1), top + Inches(0.3), card_w - Inches(0.6), Inches(0.4), "📚 KIẾN THỨC", font_name="Arial", font_size=16, font_color=CYAN_ACCENT, bold=True)
    bullets1 = [
        "Nhận biết các nguy cơ đuối nước phổ biến trong đời sống.",
        "Biết một số biện pháp phòng tránh chủ động.",
        "Biết các vật có thể hỗ trợ nổi trên mặt nước."
    ]
    add_bullet_points(slide, Inches(1.1), top + Inches(0.9), card_w - Inches(0.6), card_h - Inches(1.2), bullets1, font_size=14, font_color=WHITE)
    
    # Card 2: Skills
    card2 = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8) + card_w + spacing, top, card_w, card_h)
    card2.fill.solid()
    card2.fill.fore_color.rgb = NAVY_MEDIUM
    card2.line.color.rgb = CYAN_ACCENT
    card2.line.width = Pt(1.5)
    add_textbox(slide, Inches(1.1) + card_w + spacing, top + Inches(0.3), card_w - Inches(0.6), Inches(0.4), "🛠️ KỸ NĂNG", font_name="Arial", font_size=16, font_color=CYAN_ACCENT, bold=True)
    bullets2 = [
        "Thu thập thông tin và dữ kiện.",
        "Phân tích tình huống nguy hiểm chi tiết.",
        "Kết nối nguyên nhân và giải pháp thực tế."
    ]
    add_bullet_points(slide, Inches(1.1) + card_w + spacing, top + Inches(0.9), card_w - Inches(0.6), card_h - Inches(1.2), bullets2, font_size=14, font_color=WHITE)
    
    # Card 3: Message
    card3 = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8) + 2 * (card_w + spacing), top, card_w, card_h)
    card3.fill.solid()
    card3.fill.fore_color.rgb = CORAL_WARNING
    card3.line.color.rgb = CORAL_WARNING
    add_textbox(slide, Inches(1.1) + 2 * (card_w + spacing), top + Inches(0.3), card_w - Inches(0.6), Inches(0.4), "📢 THÔNG ĐIỆP CHÍNH", font_name="Arial", font_size=16, font_color=WHITE, bold=True)
    
    msg_text = "An toàn dưới nước bắt đầu từ sự chuẩn bị của mỗi chúng ta."
    add_textbox(slide, Inches(1.1) + 2 * (card_w + spacing), top + Inches(1.4), card_w - Inches(0.6), card_h - Inches(1.8), msg_text, font_name="Arial", font_size=20, font_color=WHITE, bold=True, align=PP_ALIGN.CENTER)

def create_end02_slide(prs, image_path):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_dark_background(slide)
    
    add_textbox(slide, Inches(0.8), Inches(0.8), Inches(11.7), Inches(0.5), "SẢN PHẨM CUỐI TIẾT (FINAL OUTPUT)", font_name="Arial", font_size=24, font_color=WHITE, bold=True)
    
    # Left Card
    card_w = Inches(5.8)
    card_h = Inches(4.8)
    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(1.6), card_w, card_h)
    card.fill.solid()
    card.fill.fore_color.rgb = NAVY_MEDIUM
    card.line.color.rgb = CYAN_ACCENT
    card.line.width = Pt(1.5)
    
    add_textbox(slide, Inches(1.2), Inches(1.9), Inches(5.0), Inches(0.4), "📖 SỔ TAY PHÓNG VIÊN PHÒNG TRÁNH ĐUỐI NƯỚC", font_name="Arial", font_size=16, font_color=CYAN_ACCENT, bold=True)
    add_textbox(slide, Inches(1.2), Inches(2.5), Inches(5.0), Inches(0.4), "Thành phần chính của sản phẩm học sinh:", font_name="Arial", font_size=14, font_color=RGBColor(200, 200, 200), bold=True)
    
    components = [
        "Các nguy cơ đuối nước phổ biến trong cuộc sống.",
        "Cách phòng tránh tương ứng với từng nguy cơ.",
        "Những vật có thể hỗ trợ nổi khi gặp nguy cơ đuối nước."
    ]
    add_bullet_points(slide, Inches(1.2), Inches(3.0), Inches(5.0), Inches(3.0), components, font_size=15, font_color=WHITE)
    
    if os.path.exists(image_path):
        slide.shapes.add_picture(image_path, Inches(7.0), Inches(1.6), width=Inches(5.5), height=Inches(4.8))
        # Overlay details onto notebook pages
        # Left Page text
        left_txt = "📓 SỔ TAY PHÒNG TRÁNH\n\n🔍 Chuột rút ➔ Khởi động\n🌪️ Nước xoáy ➔ Tránh xa\n⛵ Lật thuyền ➔ Mặc áo phao"
        add_textbox(slide, Inches(7.35), Inches(2.1), Inches(2.2), Inches(3.8), left_txt, font_name="Arial", font_size=10, font_color=DARK_TEXT, bold=True)
        # Right Page text
        right_txt = "🦺 VẬT HỖ TRỢ NỔI\n\n✔ Áo phao\n✔ Can nhựa rỗng\n✔ Khúc cây nổi\n✔ Cành cây"
        add_textbox(slide, Inches(9.85), Inches(2.1), Inches(2.2), Inches(3.8), right_txt, font_name="Arial", font_size=10, font_color=DARK_TEXT, bold=True)
    else:
        pl = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(7.0), Inches(1.6), Inches(5.5), Inches(4.8))
        pl.fill.solid()
        pl.fill.fore_color.rgb = NAVY_MEDIUM
        pl.line.color.rgb = CYAN_ACCENT
        add_textbox(slide, Inches(7.2), Inches(3.6), Inches(5.1), Inches(1.0), f"[Visual Notebook Output]\n{image_path}", font_name="Arial", font_size=16, bold=True, align=PP_ALIGN.CENTER)

def create_end03_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_dark_background(slide)
    
    add_textbox(slide, Inches(0.8), Inches(0.8), Inches(11.7), Inches(0.5), "QUAY LẠI CÂU HỎI LỚN (BIG QUESTION REVISIT)", font_name="Arial", font_size=24, font_color=WHITE, bold=True)
    
    # Big Question Box
    q_card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1.5), Inches(1.8), Inches(10.333), Inches(1.5))
    q_card.fill.solid()
    q_card.fill.fore_color.rgb = NAVY_MEDIUM
    q_card.line.color.rgb = CORAL_WARNING
    q_card.line.width = Pt(2)
    
    add_textbox(slide, Inches(1.8), Inches(2.0), Inches(9.733), Inches(0.4), "❓ CÂU HỎI LỚN ĐẦU GIỜ:", font_name="Arial", font_size=13, font_color=CORAL_WARNING, bold=True)
    add_textbox(slide, Inches(1.8), Inches(2.4), Inches(9.733), Inches(0.7), "Làm thế nào để phòng tránh đuối nước và giúp người khác an toàn hơn?", font_name="Arial", font_size=20, font_color=WHITE, bold=True)
    
    # Answer Box
    a_card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1.5), Inches(3.8), Inches(10.333), Inches(2.6))
    a_card.fill.solid()
    a_card.fill.fore_color.rgb = CYAN_ACCENT
    a_card.line.color.rgb = CYAN_ACCENT
    
    add_textbox(slide, Inches(1.8), Inches(4.0), Inches(9.733), Inches(0.4), "💡 CÂU TRẢ LỜI SAU TIẾT HỌC:", font_name="Arial", font_size=13, font_color=NAVY_DARK, bold=True)
    
    insight_text = "Đuối nước không chỉ xảy ra khi bơi. Nhiều tình huống khác nhau đều có thể dẫn đến đuối nước. Nhận biết nguy cơ sớm và chuẩn bị trước là cách tốt nhất để bảo vệ bản thân."
    add_textbox(slide, Inches(1.8), Inches(4.5), Inches(9.733), Inches(1.7), "📢 " + insight_text, font_name="Arial", font_size=22, font_color=WHITE, bold=True)

def create_end04_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_dark_background(slide)
    
    add_textbox(slide, Inches(0.8), Inches(0.8), Inches(11.7), Inches(0.5), "TỰ PHẢN CHIẾU (REFLECTION)", font_name="Arial", font_size=24, font_color=WHITE, bold=True)
    
    card_w = Inches(5.6)
    card_h = Inches(2.2)
    left_col = Inches(0.8)
    right_col = Inches(6.9)
    top_row = Inches(1.6)
    bot_row = Inches(4.2)
    
    # Card 1: Bất ngờ nhất
    c1 = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left_col, top_row, card_w, card_h)
    c1.fill.solid()
    c1.fill.fore_color.rgb = NAVY_MEDIUM
    c1.line.color.rgb = CYAN_ACCENT
    c1.line.width = Pt(1.5)
    add_textbox(slide, left_col + Inches(0.3), top_row + Inches(0.2), card_w - Inches(0.6), Inches(0.3), "⭐ ĐIỀU EM BẤT NGỜ NHẤT LÀ...", font_name="Arial", font_size=13, font_color=CYAN_ACCENT, bold=True)
    add_textbox(slide, left_col + Inches(0.3), top_row + Inches(0.6), card_w - Inches(0.6), Inches(1.4), "⭐ Điền theo phiếu Reflection của học sinh: Các góc nhìn mới mẻ mà học sinh đã thu hoạch sau tiết học.", font_name="Calibri", font_size=14, font_color=RGBColor(160, 174, 192))
    
    # Card 2: Trước đây nghĩ
    c2 = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, right_col, top_row, card_w, card_h)
    c2.fill.solid()
    c2.fill.fore_color.rgb = NAVY_MEDIUM
    c2.line.color.rgb = CYAN_ACCENT
    c2.line.width = Pt(1.5)
    add_textbox(slide, right_col + Inches(0.3), top_row + Inches(0.2), card_w - Inches(0.6), Inches(0.3), "💭 TRƯỚC ĐÂY EM NGHĨ...", font_name="Arial", font_size=13, font_color=CYAN_ACCENT, bold=True)
    add_textbox(slide, right_col + Inches(0.3), top_row + Inches(0.6), card_w - Inches(0.6), Inches(1.4), "💭 Em nghĩ đuối nước chỉ xảy ra khi đi bơi ở chỗ nước sâu, hoặc do bản thân không biết bơi. Chỉ cần đứng xa rìa nước hoặc biết bơi cơ bản là an toàn.", font_name="Calibri", font_size=14, font_color=WHITE)
    
    # Card 3: Bây giờ hiểu
    c3 = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left_col, bot_row, card_w, card_h)
    c3.fill.solid()
    c3.fill.fore_color.rgb = NAVY_MEDIUM
    c3.line.color.rgb = CYAN_ACCENT
    c3.line.width = Pt(1.5)
    add_textbox(slide, left_col + Inches(0.3), bot_row + Inches(0.2), card_w - Inches(0.6), Inches(0.3), "💡 BÂY GIỜ EM HIỂU RẰNG...", font_name="Arial", font_size=13, font_color=CYAN_ACCENT, bold=True)
    add_textbox(slide, left_col + Inches(0.3), bot_row + Inches(0.6), card_w - Inches(0.6), Inches(1.4), "💡 Nhiều nguy cơ đuối nước có thể xuất hiện trong cuộc sống hằng ngày (như chuột rút khi bơi, gặp dòng nước xoáy nguy hiểm, bị lật thuyền).", font_name="Calibri", font_size=14, font_color=WHITE)
    
    # Card 4: Role Lens
    c4 = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, right_col, bot_row, card_w, card_h)
    c4.fill.solid()
    c4.fill.fore_color.rgb = CORAL_WARNING
    c4.line.color.rgb = CORAL_WARNING
    add_textbox(slide, right_col + Inches(0.3), bot_row + Inches(0.2), card_w - Inches(0.6), Inches(0.3), "🎤 GÓC NHÌN PHÓNG VIÊN (ROLE LENS)", font_name="Arial", font_size=13, font_color=WHITE, bold=True)
    add_textbox(slide, right_col + Inches(0.3), bot_row + Inches(0.6), card_w - Inches(0.6), Inches(1.4), "🎤 Nếu là phóng viên Nova News, em muốn cảnh báo cộng đồng về nguy cơ nào nhất? Vì sao?", font_name="Calibri", font_size=14, font_color=WHITE)

def create_end05_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_dark_background(slide)
    
    add_textbox(slide, Inches(0.8), Inches(0.8), Inches(11.7), Inches(0.5), "CAM KẾT HÀNH ĐỘNG (ACTION COMMITMENT)", font_name="Arial", font_size=24, font_color=WHITE, bold=True)
    
    card_w = Inches(3.6)
    card_h = Inches(3.6)
    spacing = Inches(0.45)
    top = Inches(1.8)
    
    commitments = [
        ("👁️ CAM KẾT 1", "Quan sát kỹ lưỡng các nguy cơ khi đến gần khu vực sông, hồ, ao, suối, biển..."),
        ("🚣 CAM KẾT 2", "Luôn mặc áo phao đúng quy định và cẩn thận khi di chuyển bằng thuyền, đò."),
        ("🗣️ CAM KẾT 3", "Nhắc nhở bạn bè và người thân xung quanh thực hiện các quy tắc an toàn phòng tránh đuối nước.")
    ]
    
    for i, (title, text) in enumerate(commitments):
        left = Inches(0.8) + i * (card_w + spacing)
        card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, card_w, card_h)
        card.fill.solid()
        card.fill.fore_color.rgb = NAVY_MEDIUM
        card.line.color.rgb = CYAN_ACCENT
        card.line.width = Pt(1.5)
        
        add_textbox(slide, left + Inches(0.3), top + Inches(0.4), card_w - Inches(0.6), Inches(0.4), title, font_name="Arial", font_size=18, font_color=CYAN_ACCENT, bold=True, align=PP_ALIGN.CENTER)
        add_textbox(slide, left + Inches(0.3), top + Inches(1.2), card_w - Inches(0.6), card_h - Inches(1.6), text, font_name="Arial", font_size=16, font_color=WHITE, align=PP_ALIGN.CENTER)
        
    support_card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(5.8), Inches(11.7), Inches(0.8))
    support_card.fill.solid()
    support_card.fill.fore_color.rgb = CORAL_WARNING
    support_card.line.color.rgb = CORAL_WARNING
    add_textbox(slide, Inches(1.1), Inches(6.05), Inches(11.1), Inches(0.4), "🛠️ CÔNG CỤ HỖ TRỢ: Sổ tay phóng viên phòng tránh đuối nước.", font_name="Arial", font_size=16, font_color=WHITE, bold=True, align=PP_ALIGN.CENTER)

# -----------------------------------------------------------------------------
# Main Assembly Execution
# -----------------------------------------------------------------------------
def build_deck():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    
    # Directory paths for images
    img_dir = "images"
    welcome_img = os.path.join(img_dir, "welcome_screen_1781236319288.png")
    big_question_img = os.path.join(img_dir, "big_question_1781236338352.png")
    reporter_role_img = os.path.join(img_dir, "reporter_role_1781236354913.png")
    witness_profiles_img = os.path.join(img_dir, "witness_profiles_1781236367723.png")
    floating_objects_img = os.path.join(img_dir, "floating_objects_1781236383139.png")
    reporter_notebook_img = os.path.join(img_dir, "reporter_notebook_1781236401501.png")
    reflection_journal_img = os.path.join(img_dir, "reflection_journal_1781236420017.png")
    safety_commitment_img = os.path.join(img_dir, "safety_commitment_1781236438910.png")
    
    # New blank worksheet/grid graphics
    analysis_sheet_img = os.path.join(img_dir, "analysis_sheet.png")
    solution_table_img = os.path.join(img_dir, "solution_table.png")

    print("Building slide deck components with new design upgrades (v2)...")
    
    # -------------------------------------------------------------------------
    # PRE-LESSON (3 slides)
    # -------------------------------------------------------------------------
    create_welcome_slide(prs, welcome_img)
    create_big_question_slide(prs, big_question_img)
    create_role_intro_slide(prs, reporter_role_img)
    
    # -------------------------------------------------------------------------
    # GIAI ĐOẠN 1 – MỞ KHÓA NHIỆM VỤ (6 slides)
    # -------------------------------------------------------------------------
    create_stage_intro_slide(prs, 1, "MỞ KHÓA NHIỆM VỤ", "Situation")
    create_act01_slide(prs, 1, "MỞ KHÓA NHIỆM VỤ", "TIN NÓNG", 
                       "• Khơi gợi sự quan tâm tới chủ đề\n• Tạo nhu cầu tìm hiểu", 
                       big_question_img, 
                       "Một nhóm học sinh đi chơi gần sông. Một bạn bị trượt chân xuống nước. Các bạn còn lại không biết làm gì để cứu bạn.")
    create_act02_slide(prs, 1, "MỞ KHÓA NHIỆM VỤ", "TIN NÓNG", ["Đọc bản tin kỹ lưỡng", "Chia sẻ suy nghĩ ban đầu với cả lớp"])
    create_act03_slide(prs, 1, "MỞ KHÓA NHIỆM VỤ", "TIN NÓNG", 
                       "Thảo luận:\n- Chuyện gì đã xảy ra?\n- Nếu em ở đó, em sẽ làm gì?", 
                       "Các dự đoán ban đầu về nguyên nhân và cách xử lý", "Cả lớp", "5 phút")
    create_act04_slide(prs, 1, "MỞ KHÓA NHIỆM VỤ", "TIN NÓNG", 
                       "• Học sinh đọc bản tin\n• Học sinh chia sẻ suy nghĩ\n• Các dự đoán ban đầu về nguyên nhân và cách xử lý",
                       "Cần tìm hiểu nguyên nhân và cách xử lý để bảo vệ bản thân và người khác khi gặp nguy cơ đuối nước.")
    create_act05_slide(prs, 1, "MỞ KHÓA NHIỆM VỤ", "TIN NÓNG", 
                       "Một tình huống đuối nước có thể xảy ra trong đời sống.", 
                       "Phóng viên cần nhận nhiệm vụ điều tra điều gì?")

    # -------------------------------------------------------------------------
    # GIAI ĐOẠN 2 – NHẬN NHIỆM VỤ PHÓNG VIÊN (6 slides)
    # -------------------------------------------------------------------------
    create_stage_intro_slide(prs, 2, "NHẬN NHIỆM VỤ PHÓNG VIÊN", "Situation")
    create_act01_slide(prs, 2, "NHẬN NHIỆM VỤ PHÓNG VIÊN", "NHẬN NHIỆM VỤ TỪ TÒA SOẠN",
                       "• Hiểu chủ đề điều tra\n• Xác định câu hỏi trọng tâm",
                       welcome_img,
                       "Nhiệm vụ từ Nova News: điều tra nguy cơ đuối nước, tìm cách phòng tránh, thu thập thông tin giúp trẻ em an toàn hơn.")
    create_act02_slide(prs, 2, "NHẬN NHIỆM VỤ PHÓNG VIÊN", "NHẬN NHIỆM VỤ TỪ TÒA SOẠN", ["Tiếp nhận nhiệm vụ từ tòa soạn"])
    create_act03_slide(prs, 2, "NHẬN NHIỆM VỤ PHÓNG VIÊN", "NHẬN NHIỆM VỤ TỪ TÒA SOẠN",
                       "Nova News giao nhiệm vụ: điều tra các nguy cơ đuối nước; tìm cách phòng tránh; thu thập thông tin giúp trẻ em an toàn hơn.\n\nCâu hỏi điều tra:\n- Những nguy cơ đuối nước nào thường gặp?\n- Vì sao chúng nguy hiểm?\n- Làm thế nào để phòng tránh?",
                       "Bộ câu hỏi điều tra", "Cả lớp", "4 phút")
    create_act04_slide(prs, 2, "NHẬN NHIỆM VỤ PHÓNG VIÊN", "NHẬN NHIỆM VỤ TỪ TÒA SOẠN",
                       "• Điều tra các nguy cơ đuối nước\n• Tìm cách phòng tránh\n• Thu thập thông tin giúp trẻ em an toàn hơn",
                       "Phóng viên cần xác định đúng câu hỏi điều tra trước khi thu thập thông tin.")
    create_act05_slide(prs, 2, "NHẬN NHIỆM VỤ PHÓNG VIÊN", "NHẬN NHIỆM VỤ TỪ TÒA SOẠN",
                       "Nhiệm vụ và bộ câu hỏi điều tra của phóng viên Nova News.",
                       "Những nhân chứng đã gặp tình huống nguy hiểm nào?")

    # -------------------------------------------------------------------------
    # GIAI ĐOẠN 3 – THU THẬP THÔNG TIN (6 slides)
    # -------------------------------------------------------------------------
    create_stage_intro_slide(prs, 3, "THU THẬP THÔNG TIN", "Experience")
    create_act01_slide(prs, 3, "THU THẬP THÔNG TIN", "PHỎNG VẤN NHÂN CHỨNG",
                       "• Thu thập dữ kiện\n• Xác định các nguy cơ đuối nước",
                       witness_profiles_img,
                       "Hồ sơ nhân chứng:\n- Nhân chứng 1 bị chuột rút khi đang bơi\n- Nhân chứng 2 gặp dòng nước xoáy\n- Nhân chứng 3 thuyền bị lật")
    create_act02_slide(prs, 3, "THU THẬP THÔNG TIN", "PHỎNG VẤN NHÂN CHỨNG", ["Đọc hồ sơ nhân chứng", "Xác định các nguy cơ cụ thể"])
    create_act03_slide(prs, 3, "THU THẬP THÔNG TIN", "PHỎNG VẤN NHÂN CHỨNG",
                       "Mỗi nhóm nhận hồ sơ nhân chứng. Nhiệm vụ:\n- Nhân chứng đã gặp tình huống nào?\n- Em phát hiện những nguy cơ đuối nước nào?\n\nCác nhóm ghi kết quả lên phiếu điều tra.",
                       "Danh sách nguy cơ đuối nước", "Nhóm 4", "7 phút")
    create_act04_slide(prs, 3, "THU THẬP THÔNG TIN", "PHỎNG VẤN NHÂN CHỨNG",
                       "• Chuột rút khi đang bơi\n• Dòng nước xoáy\n• Thuyền bị lật",
                       "Thu thập dữ kiện giúp phóng viên nhận diện các nguy cơ trước khi phân tích.")
    create_act05_slide(prs, 3, "THU THẬP THÔNG TIN", "PHỎNG VẤN NHÂN CHỨNG",
                       "Các tình huống nhân chứng gặp phải và danh sách nguy cơ đuối nước.",
                       "Vì sao những tình huống này nguy hiểm?")

    # -------------------------------------------------------------------------
    # GIAI ĐOẠN 4 – KHÁM PHÁ CÁC GÓC NHÌN (6 slides - Overlaying native table)
    # -------------------------------------------------------------------------
    create_stage_intro_slide(prs, 4, "KHÁM PHÁ CÁC GÓC NHÌN", "Thinking")
    
    # Overlaid Table Data for S4-ACT-01
    s4_headers = ["Nguy cơ", "Điều có thể xảy ra"]
    s4_rows = [
        ["Chuột rút khi đang bơi", "Cơ co cứng đau đớn, mất khả năng tự nổi và chìm nhanh."],
        ["Gặp dòng nước xoáy", "Bị cuốn xuống đáy sâu, cản trở việc ngoi lên thở."],
        ["Thuyền bị lật bất ngờ", "Rơi bất ngờ, hoảng loạn, không điểm tựa và chìm."]
    ]
    
    create_act01_slide(prs, 4, "KHÁM PHÁ CÁC GÓC NHÌN", "NHÌN TỪ GÓC ĐỘ NHÂN CHỨNG",
                       "• Hiểu vì sao các tình huống đó nguy hiểm\n• Nhìn vấn đề từ góc độ người gặp nạn",
                       analysis_sheet_img,
                       "Phiếu phân tích nguy cơ và hậu quả.",
                       table_info=(s4_headers, s4_rows))
                       
    create_act02_slide(prs, 4, "KHÁM PHÁ CÁC GÓC NHÌN", "NHÌN TỪ GÓC ĐỘ NHÂN CHỨNG", ["Phân tích hậu quả nguy hiểm của từng trường hợp"])
    create_act03_slide(prs, 4, "KHÁM PHÁ CÁC GÓC NHÌN", "NHÌN TỪ GÓC ĐỘ NHÂN CHỨNG",
                       "Thảo luận:\n- Vì sao chuột rút khi đang bơi lại nguy hiểm?\n- Điều gì có thể xảy ra khi gặp nước xoáy?\n- Người trên thuyền bị lật sẽ gặp khó khăn gì?\n\nHoàn thành bảng Nguy cơ ➔ Điều có thể xảy ra.",
                       "Bảng phân tích nguy cơ và hậu quả", "Nhóm 4", "6 phút")
    create_act04_slide(prs, 4, "KHÁM PHÁ CÁC GÓC NHÌN", "NHÌN TỪ GÓC ĐỘ NHÂN CHỨNG",
                       "• Chuột rút ➔ Điều có thể xảy ra: Mất khả năng bơi và nổi.\n• Nước xoáy ➔ Điều có thể xảy ra: Cuốn cơ thể chìm xuống đáy.\n• Lật thuyền ➔ Điều có thể xảy ra: Hoảng loạn, uống nước và ngạt thở.",
                       "Những tình huống này đều có thể khiến người gặp nạn mất khả năng tự bảo vệ bản thân dưới nước.")
    create_act05_slide(prs, 4, "KHÁM PHÁ CÁC GÓC NHÌN", "NHÌN TỪ GÓC ĐỘ NHÂN CHỨNG",
                       "Hậu quả có thể xảy ra từ từng nguy cơ đuối nước.",
                       "Từ các nguy cơ vừa phân tích, chúng ta nên làm gì để phòng tránh?")

    # -------------------------------------------------------------------------
    # GIAI ĐOẠN 5 – TÌM RA ĐIỀU QUAN TRỌNG (11 slides - 2 Activities)
    # -------------------------------------------------------------------------
    create_stage_intro_slide(prs, 5, "TÌM RA ĐIỀU QUAN TRỌNG", "Decision")
    
    # Act 5.1: Tìm giải pháp (Overlaying native table)
    s5_headers = ["Nguy cơ", "Cách phòng tránh"]
    s5_rows = [
        ["Bị chuột rút", "Khởi động kỹ và làm ấm cơ thể trước khi xuống nước."],
        ["Dòng nước xoáy", "Tránh xa và tuyệt đối không tắm ở khu vực nguy hiểm."],
        ["Thuyền bị lật", "Luôn mặc áo phao đúng cách suốt chuyến đi."]
    ]
    
    create_act01_slide(prs, 5, "TÌM RA ĐIỀU QUAN TRỌNG", "TÌM GIẢI PHÁP",
                       "• Xác định thông tin quan trọng nhất\n• Tìm giải pháp phòng tránh hiệu quả",
                       solution_table_img,
                       "Bảng kết nối Nguy cơ ➔ Cách phòng tránh.",
                       table_info=(s5_headers, s5_rows))
                       
    create_act02_slide(prs, 5, "TÌM RA ĐIỀU QUAN TRỌNG", "TÌM GIẢI PHÁP", ["Nhóm thảo luận: Từ các nguy cơ đã biết, làm sao để phòng tránh?"])
    create_act03_slide(prs, 5, "TÌM RA ĐIỀU QUAN TRỌNG", "TÌM GIẢI PHÁP",
                       "Từ các nguy cơ vừa phân tích, chúng ta nên làm gì để phòng tránh?\n\nHoàn thành bảng kết nối Nguy cơ ➔ Cách phòng tránh.",
                       "Danh sách giải pháp phòng tránh", "Nhóm", "6 phút")
    create_act04_slide(prs, 5, "TÌM RA ĐIỀU QUAN TRỌNG", "TÌM GIẢI PHÁP",
                       "• Khởi động trước khi xuống nước\n• Mặc áo phao khi đi thuyền\n• Không chơi gần khu vực nguy hiểm",
                       "Kết nối nguy cơ với giải pháp giúp phòng tránh đuối nước chủ động hơn.")
    create_act05_slide(prs, 5, "TÌM RA ĐIỀU QUAN TRỌNG", "TÌM GIẢI PHÁP",
                       "Các giải pháp phòng tránh tương ứng với nguy cơ.",
                       "Những vật nào có thể hỗ trợ nổi trên mặt nước?")

    # Act 5.2: Vật nào giúp mình nổi?
    create_act01_slide(prs, 5, "TÌM RA ĐIỀU QUAN TRỌNG", "VẬT NÀO GIÚP MÌNH NỔI?",
                       "• Xác định thông tin quan trọng nhất\n• Tìm giải pháp phòng tránh",
                       floating_objects_img,
                       "Thẻ hình ảnh: áo phao, can nhựa rỗng, khúc cây nổi, cục đá, cành cây chìa ra bờ.")
    create_act02_slide(prs, 5, "TÌM RA ĐIỀU QUAN TRỌNG", "VẬT NÀO GIÚP MÌNH NỔI?", ["Phân loại các vật vào 2 nhóm: Giúp nổi / Không giúp nổi"])
    create_act03_slide(prs, 5, "TÌM RA ĐIỀU QUAN TRỌNG", "VẬT NÀO GIÚP MÌNH NỔI?",
                       "Phân loại các vật: áo phao, can nhựa rỗng, khúc cây nổi, cục đá, cành cây chìa ra bờ vào 2 nhóm: Giúp nổi / Không giúp nổi.",
                       "Danh sách vật hỗ trợ nổi", "Nhóm", "6 phút")
    create_act04_slide(prs, 5, "TÌM RA ĐIỀU QUAN TRỌNG", "VẬT NÀO GIÚP MÌNH NỔI?",
                       "• Giúp nổi: áo phao, can nhựa rỗng, khúc cây nổi, cành cây chìa ra bờ\n• Không giúp nổi: cục đá\n• Điểm chung: là vật nổi hoặc giúp cơ thể duy trì trên mặt nước",
                       "Vật hỗ trợ cần là vật nổi hoặc giúp cơ thể duy trì trên mặt nước.")
    create_act05_slide(prs, 5, "TÌM RA ĐIỀU QUAN TRỌNG", "VẬT NÀO GIÚP MÌNH NỔI?",
                       "Các vật có thể và không thể hỗ trợ nổi.",
                       "Làm thế nào chuyển thông tin đã điều tra thành sổ tay hữu ích?")

    # -------------------------------------------------------------------------
    # GIAI ĐOẠN 6 – THÔNG ĐIỆP PHÓNG VIÊN (6 slides - Overlaying text boxes on notebook pages)
    # -------------------------------------------------------------------------
    create_stage_intro_slide(prs, 6, "THÔNG ĐIỆP PHÓNG VIÊN", "Action")
    
    s6_left_page = "📓 SỔ TAY PHÒNG TRÁNH\n\n🔍 Chuột rút ➔ Khởi động\n🌪️ Nước xoáy ➔ Tránh xa\n⛵ Lật thuyền ➔ Mặc áo phao"
    s6_right_page = "🦺 VẬT HỖ TRỢ NỔI\n\n✔ Áo phao\n✔ Can nhựa rỗng\n✔ Khúc cây nổi\n✔ Cành cây"
    
    create_act01_slide(prs, 6, "THÔNG ĐIỆP PHÓNG VIÊN", "SỔ TAY PHÓNG VIÊN",
                       "• Chuyển thông tin thành sản phẩm hữu ích",
                       reporter_notebook_img,
                       "Bố cục sổ tay gồm Nguy cơ ➔ Cách phòng tránh và Vật hỗ trợ nổi.",
                       notebook_info=(s6_left_page, s6_right_page))
                       
    create_act02_slide(prs, 6, "THÔNG ĐIỆP PHÓNG VIÊN", "SỔ TAY PHÓNG VIÊN", ["Nhóm hoàn thiện sổ tay", "Đại diện nhóm chia sẻ"])
    create_act03_slide(prs, 6, "THÔNG ĐIỆP PHÓNG VIÊN", "SỔ TAY PHÓNG VIÊN",
                       "Nhóm hoàn thiện sổ tay: Chuột rút ➔ Khởi động trước khi xuống nước; Nước xoáy ➔ Tránh khu vực nguy hiểm; Lật thuyền ➔ Mặc áo phao.\n\nVật hỗ trợ nổi: áo phao, can rỗng, khúc cây, cành cây.\n\nĐại diện nhóm chia sẻ trước lớp.",
                       "Sổ tay phóng viên hoàn chỉnh", "Nhóm 4", "10 phút")
    create_act04_slide(prs, 6, "THÔNG ĐIỆP PHÓNG VIÊN", "SỔ TAY PHÓNG VIÊN",
                       "• Nguy cơ ➔ Cách phòng tránh\n• Vật hỗ trợ nổi\n• Đại diện nhóm chia sẻ sổ tay",
                       "Thông tin điều tra cần được chuyển thành sản phẩm dễ hiểu, giúp cảnh báo và hỗ trợ người khác.")
    create_act05_slide(prs, 6, "THÔNG ĐIỆP PHÓNG VIÊN", "SỔ TAY PHÓNG VIÊN",
                       "Sổ tay phóng viên phòng tránh đuối nước.",
                       "Điều gì khiến em bất ngờ và muốn ghi nhớ sau bài học?")

    # -------------------------------------------------------------------------
    # GIAI ĐOẠN 7 – NHẬT KÝ PHÓNG VIÊN (6 slides - Overlaying text box on journal page)
    # -------------------------------------------------------------------------
    create_stage_intro_slide(prs, 7, "NHẬT KÝ PHÓNG VIÊN", "Reflection")
    
    s7_journal_text = "📓 NHẬT KÝ PHÓNG VIÊN\n\n⭐ Điều làm em bất ngờ nhất là...\n\n⭐ Nguy cơ em chưa từng nghĩ tới là...\n\n⭐ Điều em muốn ghi nhớ là..."
    
    create_act01_slide(prs, 7, "NHẬT KÝ PHÓNG VIÊN", "NHẬT KÝ PHÓNG VIÊN",
                       "• Thực hiện hoạt động tự phản chiếu (Reflection)",
                       reflection_journal_img,
                       "Phiếu Reflection theo dạng nhật ký phóng viên.",
                       journal_info=s7_journal_text)
                       
    create_act02_slide(prs, 7, "NHẬT KÝ PHÓNG VIÊN", "NHẬT KÝ PHÓNG VIÊN", 
                       ["Hoàn thành: Điều làm em bất ngờ nhất là...", 
                        "Hoàn thành: Nguy cơ em chưa từng nghĩ tới là...", 
                        "Hoàn thành: Điều em muốn ghi nhớ là..."])
    create_act03_slide(prs, 7, "NHẬT KÝ PHÓNG VIÊN", "NHẬT KÝ PHÓNG VIÊN",
                       "Hoàn thành phiếu Reflection: Điều làm em bất ngờ nhất là...; Nguy cơ em chưa từng nghĩ tới là...; Điều em muốn ghi nhớ là...",
                       "Phiếu Reflection hoàn thiện", "Cá nhân", "4 phút")
    create_act04_slide(prs, 7, "NHẬT KÝ PHÓNG VIÊN", "NHẬT KÝ PHÓNG VIÊN",
                       "• Điều làm em bất ngờ nhất\n• Nguy cơ em chưa từng nghĩ tới\n• Điều em muốn ghi nhớ",
                       "Reflection giúp học sinh nhận ra nhiều nguy cơ đuối nước có thể xuất hiện trong cuộc sống hằng ngày.")
    create_act05_slide(prs, 7, "NHẬT KÝ PHÓNG VIÊN", "NHẬT KÝ PHÓNG VIÊN",
                       "Những điều bất ngờ, nguy cơ chưa từng nghĩ tới và điều cần ghi nhớ.",
                       "Sau bài học này, em sẽ hành động như thế nào?")

    # -------------------------------------------------------------------------
    # GIAI ĐOẠN 8 – ÁP DỤNG NGAY (6 slides - Overlaying text box on commitment card)
    # -------------------------------------------------------------------------
    create_stage_intro_slide(prs, 8, "ÁP DỤNG NGAY", "Adaptation")
    
    s8_commitment_text = "📜 BẢN CAM KẾT AN TOÀN\n\n👁️ Em sẽ...\n\n👁️ Em sẽ nhắc bạn...\n\n👁️ Em sẽ báo người lớn khi..."
    
    create_act01_slide(prs, 8, "ÁP DỤNG NGAY", "CAM KẾT AN TOÀN",
                       "• Chuyển bài học thành hành động thực tế",
                       safety_commitment_img,
                       "Phiếu cam kết an toàn cá nhân.",
                       commitment_info=s8_commitment_text)
                       
    create_act02_slide(prs, 8, "ÁP DỤNG NGAY", "CAM KẾT AN TOÀN", 
                       ["Hoàn thành: Em sẽ...", 
                        "Hoàn thành: Em sẽ nhắc bạn...", 
                        "Hoàn thành: Em sẽ báo người lớn khi..."])
    create_act03_slide(prs, 8, "ÁP DỤNG NGAY", "CAM KẾT AN TOÀN",
                       "Hoàn thành phiếu cam kết: Sau bài học này: Em sẽ...; Em sẽ nhắc bạn...; Em sẽ báo người lớn khi...",
                       "Cam kết hành động", "Cá nhân", "3 phút")
    create_act04_slide(prs, 8, "ÁP DỤNG NGAY", "CAM KẾT AN TOÀN",
                       "• Em sẽ...\n• Em sẽ nhắc bạn...\n• Em sẽ báo người lớn khi...",
                       "An toàn dưới nước bắt đầu từ sự chuẩn bị của mỗi chúng ta.")
    create_act05_slide(prs, 8, "ÁP DỤNG NGAY", "CAM KẾT AN TOÀN",
                       "Cam kết hành động sau bài học.",
                       "Kết thúc bài học.")

    # -------------------------------------------------------------------------
    # END LESSON (5 slides)
    # -------------------------------------------------------------------------
    create_end01_slide(prs)
    create_end02_slide(prs, reporter_notebook_img)
    create_end03_slide(prs)
    create_end04_slide(prs)
    create_end05_slide(prs)

    # -------------------------------------------------------------------------
    # Save Presentation
    # -------------------------------------------------------------------------
    output_filename = "Slide_Phong_tranh_duoi_nuoc_Phong_vien.pptx"
    prs.save(output_filename)
    print(f"Presentation updated and saved successfully as '{output_filename}'")
    print(f"Total slides generated: {len(prs.slides)}")

if __name__ == "__main__":
    build_deck()
